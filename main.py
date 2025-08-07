import gradio as gr
import pandas as pd
import matplotlib.pyplot as plt
import seaborn as sns
from openai import OpenAI
import io
import base64
from reportlab.lib.pagesizes import letter, A4
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Image, PageBreak
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.lib import colors
import docx
from pptx import Presentation
import json
import tempfile
import os
from typing import List, Dict, Any
import numpy as np
import plotly.graph_objects as go
import plotly.express as px
from plotly.offline import plot
import warnings
warnings.filterwarnings('ignore')

# Get OpenAI API key from environment
try:
    OPENAI_API_KEY = os.getenv('OPENAI_API_KEY')
    print("✅ OpenAI API key loaded")
except Exception as e:
    print("❌ Error loading OpenAI API key:", e)
    OPENAI_API_KEY = None

class ReportGenerator:
    def __init__(self, api_key: str):
        self.client = OpenAI(api_key=api_key)
        self.styles = getSampleStyleSheet()
        self.charts = []

    def extract_text_from_files(self, files: List[str]) -> Dict[str, str]:
        """Extract text content from various file types"""
        content = {}

        for file_path in files:
            file_name = os.path.basename(file_path)
            file_ext = os.path.splitext(file_path)[1].lower()

            try:
                if file_ext == '.txt':
                    with open(file_path, 'r', encoding='utf-8') as f:
                        content[file_name] = f.read()

                elif file_ext in ['.docx', '.doc']:
                    doc = docx.Document(file_path)
                    text = '\n'.join([paragraph.text for paragraph in doc.paragraphs])
                    content[file_name] = text

                elif file_ext in ['.xlsx', '.xls', '.csv']:
                    if file_ext == '.csv':
                        df = pd.read_csv(file_path)
                    else:
                        df = pd.read_excel(file_path)

                    # Convert DataFrame to readable format
                    text_content = f"Dataset Summary:\n"
                    text_content += f"Shape: {df.shape}\n"
                    text_content += f"Columns: {list(df.columns)}\n\n"
                    text_content += "First 10 rows:\n"
                    text_content += df.head(10).to_string()

                    if df.select_dtypes(include=[np.number]).shape[1] > 0:
                        text_content += "\n\nNumerical Summary:\n"
                        text_content += df.describe().to_string()

                    content[file_name] = text_content

                elif file_ext in ['.pptx', '.ppt']:
                    prs = Presentation(file_path)
                    text = ""
                    for slide in prs.slides:
                        for shape in slide.shapes:
                            if hasattr(shape, "text"):
                                text += shape.text + "\n"
                    content[file_name] = text

            except Exception as e:
                content[file_name] = f"Error reading file: {str(e)}"

        return content

    def generate_charts_from_data(self, files: List[str]) -> List[str]:
        """Generate diverse charts from data files and return base64 encoded images"""
        chart_images = []

        for file_path in files:
            file_ext = os.path.splitext(file_path)[1].lower()

            if file_ext in ['.xlsx', '.xls', '.csv']:
                try:
                    if file_ext == '.csv':
                        df = pd.read_csv(file_path)
                    else:
                        df = pd.read_excel(file_path)

                    numeric_cols = df.select_dtypes(include=[np.number]).columns
                    categorical_cols = df.select_dtypes(include=['object']).columns

                    # Chart 1: Correlation heatmap (if sufficient numeric data)
                    if len(numeric_cols) >= 2:
                        plt.figure(figsize=(10, 8))
                        correlation_matrix = df[numeric_cols].corr()
                        mask = np.triu(np.ones_like(correlation_matrix, dtype=bool))
                        sns.heatmap(correlation_matrix, mask=mask, annot=True, cmap='RdBu_r', center=0,
                                  square=True, linewidths=0.5, cbar_kws={"shrink": .8})
                        plt.title('Correlation Matrix Heatmap', fontsize=16, fontweight='bold')

                        buffer = io.BytesIO()
                        plt.savefig(buffer, format='png', dpi=300, bbox_inches='tight')
                        buffer.seek(0)
                        chart_base64 = base64.b64encode(buffer.getvalue()).decode()
                        chart_images.append(chart_base64)
                        plt.close()

                    # Chart 2: Multi-variable distribution plot
                    if len(numeric_cols) >= 1:
                        plt.figure(figsize=(15, 10))
                        n_cols = min(4, len(numeric_cols))
                        n_rows = (len(numeric_cols[:4]) + 1) // 2

                        for i, col in enumerate(numeric_cols[:4]):
                            plt.subplot(n_rows, 2, i+1)
                            # Box plot with violin overlay
                            sns.violinplot(y=df[col].dropna(), color='skyblue', alpha=0.6)
                            sns.boxplot(y=df[col].dropna(), width=0.3, boxprops=dict(alpha=0.7))
                            plt.title(f'Distribution Analysis: {col}', fontweight='bold')
                            plt.ylabel(col)

                        plt.suptitle('Statistical Distribution Analysis', fontsize=16, fontweight='bold')
                        plt.tight_layout()
                        buffer = io.BytesIO()
                        plt.savefig(buffer, format='png', dpi=300, bbox_inches='tight')
                        buffer.seek(0)
                        chart_base64 = base64.b64encode(buffer.getvalue()).decode()
                        chart_images.append(chart_base64)
                        plt.close()

                    # Chart 3: Trend analysis (time series or scatter)
                    if len(numeric_cols) >= 2:
                        plt.figure(figsize=(12, 8))
                        # Create scatter plot matrix for top numeric columns
                        cols_to_plot = numeric_cols[:3]
                        if len(cols_to_plot) >= 2:
                            # Scatter plot with trend line
                            plt.subplot(2, 1, 1)
                            x, y = cols_to_plot[0], cols_to_plot[1]
                            plt.scatter(df[x], df[y], alpha=0.6, c='coral', s=50)
                            # Add trend line
                            z = np.polyfit(df[x].dropna(), df[y].dropna(), 1)
                            p = np.poly1d(z)
                            plt.plot(df[x].sort_values(), p(df[x].sort_values()), "r--", alpha=0.8, linewidth=2)
                            plt.xlabel(x)
                            plt.ylabel(y)
                            plt.title(f'Trend Analysis: {x} vs {y}', fontweight='bold')

                            # Line plot if there's a third numeric column
                            if len(cols_to_plot) >= 3:
                                plt.subplot(2, 1, 2)
                                plt.plot(df.index, df[cols_to_plot[2]], marker='o', linewidth=2, markersize=4)
                                plt.title(f'Trend Over Time: {cols_to_plot[2]}', fontweight='bold')
                                plt.xlabel('Index')
                                plt.ylabel(cols_to_plot[2])
                                plt.grid(True, alpha=0.3)

                        plt.tight_layout()
                        buffer = io.BytesIO()
                        plt.savefig(buffer, format='png', dpi=300, bbox_inches='tight')
                        buffer.seek(0)
                        chart_base64 = base64.b64encode(buffer.getvalue()).decode()
                        chart_images.append(chart_base64)
                        plt.close()

                    # Chart 4: Categorical analysis (if categorical data exists)
                    if len(categorical_cols) >= 1 and len(df) > 1:
                        plt.figure(figsize=(14, 10))

                        # Top categories analysis
                        cat_col = categorical_cols[0]
                        top_categories = df[cat_col].value_counts().head(10)

                        plt.subplot(2, 2, 1)
                        # Horizontal bar chart
                        plt.barh(range(len(top_categories)), top_categories.values, color='lightgreen')
                        plt.yticks(range(len(top_categories)), top_categories.index)
                        plt.title(f'Top Categories: {cat_col}', fontweight='bold')
                        plt.xlabel('Count')

                        plt.subplot(2, 2, 2)
                        # Pie chart for top 6 categories
                        plt.pie(top_categories.head(6).values, labels=top_categories.head(6).index,
                               autopct='%1.1f%%', startangle=90)
                        plt.title(f'Distribution: {cat_col}', fontweight='bold')

                        # If we have both categorical and numerical data
                        if len(numeric_cols) >= 1:
                            plt.subplot(2, 1, 2)
                            # Box plot of numeric variable by category
                            num_col = numeric_cols[0]
                            # Limit to top 8 categories for clarity
                            top_cats = df[cat_col].value_counts().head(8).index
                            filtered_df = df[df[cat_col].isin(top_cats)]
                            sns.boxplot(data=filtered_df, x=cat_col, y=num_col)
                            plt.xticks(rotation=45)
                            plt.title(f'{num_col} by {cat_col}', fontweight='bold')

                        plt.tight_layout()
                        buffer = io.BytesIO()
                        plt.savefig(buffer, format='png', dpi=300, bbox_inches='tight')
                        buffer.seek(0)
                        chart_base64 = base64.b64encode(buffer.getvalue()).decode()
                        chart_images.append(chart_base64)
                        plt.close()

                    # Chart 5: Statistical summary visualization
                    if len(numeric_cols) >= 2:
                        plt.figure(figsize=(12, 8))

                        # Create a comprehensive statistical overview
                        plt.subplot(2, 2, 1)
                        # Mean values bar chart
                        means = df[numeric_cols].mean()
                        plt.bar(range(len(means)), means.values, color='gold')
                        plt.xticks(range(len(means)), means.index, rotation=45)
                        plt.title('Mean Values Comparison', fontweight='bold')
                        plt.ylabel('Mean Value')

                        plt.subplot(2, 2, 2)
                        # Standard deviation
                        stds = df[numeric_cols].std()
                        plt.bar(range(len(stds)), stds.values, color='orange')
                        plt.xticks(range(len(stds)), stds.index, rotation=45)
                        plt.title('Standard Deviation', fontweight='bold')
                        plt.ylabel('Std Dev')

                        plt.subplot(2, 1, 2)
                        # Multi-line plot for comparison
                        for col in numeric_cols[:4]:
                            normalized_data = (df[col] - df[col].min()) / (df[col].max() - df[col].min())
                            plt.plot(normalized_data.rolling(window=min(10, len(df)//5)).mean(),
                                   label=col, linewidth=2, alpha=0.8)
                        plt.title('Normalized Trends Comparison', fontweight='bold')
                        plt.xlabel('Index')
                        plt.ylabel('Normalized Value (0-1)')
                        plt.legend()
                        plt.grid(True, alpha=0.3)

                        plt.tight_layout()
                        buffer = io.BytesIO()
                        plt.savefig(buffer, format='png', dpi=300, bbox_inches='tight')
                        buffer.seek(0)
                        chart_base64 = base64.b64encode(buffer.getvalue()).decode()
                        chart_images.append(chart_base64)
                        plt.close()

                except Exception as e:
                    print(f"Error generating charts: {e}")

        return chart_images


    def analyze_with_openai(self, content: Dict[str, str], user_prompt: str, report_size: str = "standard") -> str:
        """Use OpenAI to analyze content and generate report with selected size"""

        # Prepare content for analysis
        combined_content = ""
        total_content_length = 0

        for filename, file_content in content.items():
            # Increase content preview based on report size
            content_limit = {
                "brief": 2000,
                "standard": 5000,
                "comprehensive": 8000,
                "detailed": 12000
            }.get(report_size, 5000)

            content_preview = file_content[:content_limit]
            combined_content += f"\n\n=== {filename} ===\n{content_preview}"
            total_content_length += len(file_content)

        # Report configuration based on selected size
        report_configs = {
            "brief": {
                "max_tokens": 1200,
                "sections": 4,
                "depth": "high-level",
                "style": "executive summary"
            },
            "standard": {
                "max_tokens": 2800,
                "sections": 6,
                "depth": "balanced",
                "style": "professional analysis"
            },
            "comprehensive": {
                "max_tokens": 4200,
                "sections": 8,
                "depth": "detailed",
                "style": "thorough investigation"
            },
            "detailed": {
                "max_tokens": 6000,
                "sections": 10,
                "depth": "extensive",
                "style": "in-depth research"
            }
        }

        config = report_configs.get(report_size, report_configs["standard"])

        # Adaptive system prompts based on report size
        system_prompts = {
            "brief": """You are a senior executive consultant. Generate a concise, high-impact executive report based entirely on the provided dataset and its domain.

Your primary goal is to extract **critical insights**, identify **contextual relevance**, and propose **strategic recommendations** in a structured format.

- **Dynamically create appropriate section titles** depending on the data's domain (e.g., "Placement Overview", "Revenue Trends", "User Engagement Patterns", etc.)
- If applicable, highlight statistical observations, but always tie them back to the **real-world implications** and **strategic importance**.
- Avoid generic section names like "Executive Summary" unless clearly warranted by the data.

Use **Markdown** formatting:
- Use **bold** for metrics and key terms.
- Use bullet points for actionable insights and findings.
- Keep the report short, focused, and high-impact.
""",

"standard": """You are a professional business analyst. Generate a comprehensive and insightful analytical report tailored to the content and theme of the uploaded dataset.

**Do not use fixed section titles.** Instead, analyze the file and infer relevant topic-aligned section headers.

Your report should include:
- A **concise contextual summary**
- **Key trends and insights**, not just statistics
- **Observations explained in plain language**
- Any **strategic or operational implications** discovered in the data

If appropriate, mention statistical summaries (mean, variance, correlations), but do not let them dominate the report. The goal is to **convey a story** backed by the data.

Use **Markdown formatting**:
- Section titles as headers
- Bullets for points
- Bold for important figures and terms
""",


"comprehensive": """You are a senior data scientist and strategic advisor. Generate an extensive, context-aware report from the provided dataset that demonstrates deep understanding of both the **data** and its **real-world context**.

Avoid using rigid, predefined sections. Instead:
- Identify **major content areas** and build your own section titles dynamically
- Cover **statistical, contextual, and strategic layers** of interpretation
- Highlight **dependencies, patterns, and outliers**
- Relate your analysis to **market or domain realities**, if possible

Include, but do not limit to:
- Statistical analyses and data trends
- Narrative commentary on what the data implies
- Clear **calls-to-action** or strategy suggestions
- **Time-based framing** if growth, timelines, or forecasts are relevant

Use rich **Markdown** formatting:
- ### for section headers
- **Bold** for stats and insights
- Bullet points and sublists to organize recommendations
""",

"detailed": """You are a principal research analyst and strategy architect. Generate a highly detailed report that interprets the dataset through multiple dimensions: quantitative, qualitative, contextual, and strategic.

**Do not stick to a fixed template.** Instead:
- Let the **nature of the dataset** shape your section headers and structure
- Cover everything from data quality, thematic findings, domain relevance, to advanced analytics and strategy
- Go beyond statistics; explore **what the data *means***, not just what it shows

The report should include:
- **Custom-titled sections** based on emerging patterns in the data
- If applicable, detailed statistical methods and validation
- **Comparisons**, **exceptions**, **trends**, and **predictive insights**
- **Strategic plans** with time horizons and risk outlooks

Use clear and structured **Markdown formatting** throughout:
- Use heading levels (`#`, `##`, `###`) consistently
- Bold for metrics and insights
- Sub-bullets, enumerations, and tables where useful
"""

            }

        system_prompt = system_prompts.get(report_size, system_prompts["standard"])

        user_message = f"""
        Please analyze the following content and generate a {report_size} report with {config['depth']} analysis:

        **User Requirements:** {user_prompt}

        **Content Analysis:**
        {combined_content}

        **Report Specifications:**
        - Report Type: {config['style']}
        - Target Sections: {config['sections']}
        - Analysis Depth: {config['depth']}
        - Use extensive markdown formatting with multiple header levels (# ## ### ####)
        - Include **bold** text for critical points and metrics
        - Use comprehensive bullet points and numbered lists
        - Add *emphasis* for important insights and recommendations
        - Provide specific data references and statistical observations
        - Include detailed explanations and context for each finding
        - Structure content with clear logical flow and comprehensive coverage
        - Ensure each section has substantial content (minimum 2-3 paragraphs per section)
        - Add sub-sections where appropriate for better organization

        IMPORTANT: Generate substantial content for each section. Each major section should contain multiple paragraphs with detailed analysis, specific examples, and comprehensive insights. The report should be thorough enough to span 3-4 pages when printed.
        """

        try:
            response = self.client.chat.completions.create(
                model="gpt-4",  # Using GPT-4 for better quality and longer responses
                messages=[
                    {"role": "system", "content": system_prompt},
                    {"role": "user", "content": user_message}
                ],
                max_tokens=config['max_tokens'],
                temperature=0.7
            )

            return response.choices[0].message.content

        except Exception as e:
            # Fallback to GPT-3.5-turbo if GPT-4 fails
            try:
                response = self.client.chat.completions.create(
                    model="gpt-3.5-turbo",
                    messages=[
                        {"role": "system", "content": system_prompt},
                        {"role": "user", "content": user_message}
                    ],
                    max_tokens=config['max_tokens'],
                    temperature=0.7
                )
                return response.choices[0].message.content
            except Exception as e2:
                return f"Error generating analysis: {str(e2)}\n\nPlease check your OpenAI API key and try again."



    from reportlab.platypus import Table, TableStyle

    def create_pdf_report(self, report_text: str, chart_images: List[str], title: str = "Analysis Report") -> str:
        """Create PDF report with text, charts, and embedded tables"""

        temp_file = tempfile.NamedTemporaryFile(delete=False, suffix='.pdf')
        temp_path = temp_file.name
        temp_file.close()

        doc = SimpleDocTemplate(
            temp_path,
            pagesize=A4,
            topMargin=0.75 * inch,
            bottomMargin=0.75 * inch,
            leftMargin=0.75 * inch,
            rightMargin=0.75 * inch
        )
        story = []

        custom_styles = getSampleStyleSheet()
        title_style = ParagraphStyle('CompactTitle', parent=custom_styles['Heading1'], fontSize=20,
                                    spaceAfter=18, textColor=colors.darkblue, alignment=1)
        heading1_style = ParagraphStyle('CompactHeading1', parent=custom_styles['Heading1'], fontSize=16,
                                        spaceAfter=12, spaceBefore=18, textColor=colors.darkblue)
        heading2_style = ParagraphStyle('CompactHeading2', parent=custom_styles['Heading2'], fontSize=14,
                                        spaceAfter=10, spaceBefore=14, textColor=colors.darkred)
        heading3_style = ParagraphStyle('CompactHeading3', parent=custom_styles['Heading3'], fontSize=12,
                                        spaceAfter=8, spaceBefore=12, textColor=colors.darkgreen)
        normal_style = ParagraphStyle('CompactNormal', parent=custom_styles['Normal'], fontSize=10,
                                      spaceAfter=8, leading=12)

        story.append(Paragraph(title, title_style))
        story.append(Spacer(1, 15))

        lines = report_text.split('\n')
        chart_inserted = 0
        total_charts = len(chart_images)
        table_pending = []

        def insert_table(data):
            if not data or not any(data):
                return
            table = Table(data, hAlign='LEFT')
            table.setStyle(TableStyle([
                ('BACKGROUND', (0, 0), (-1, 0), colors.lightblue),
                ('TEXTCOLOR', (0, 0), (-1, 0), colors.black),
                ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
                ('GRID', (0, 0), (-1, -1), 0.5, colors.grey),
                ('FONTSIZE', (0, 0), (-1, -1), 9),
                ('BOTTOMPADDING', (0, 0), (-1, 0), 6),
            ]))
            story.append(Spacer(1, 10))
            story.append(table)
            story.append(Spacer(1, 10))

        # Process lines and detect markdown tables
        for i, line in enumerate(lines):
            line = line.strip()
            if not line:
                continue

            if '|' in line and line.count('|') >= 2:
                # Detect markdown table row
                table_pending.append([cell.strip() for cell in line.strip('|').split('|')])
                continue
            elif table_pending:
                # End of markdown table block
                insert_table(table_pending)
                table_pending = []

            # Markdown headers
            if line.startswith('####'):
                story.append(Paragraph(line.replace('####', '').strip(), heading3_style))
            elif line.startswith('###'):
                story.append(Paragraph(line.replace('###', '').strip(), heading3_style))
            elif line.startswith('##'):
                story.append(Paragraph(line.replace('##', '').strip(), heading2_style))
            elif line.startswith('#'):
                story.append(Paragraph(line.replace('#', '').strip(), heading1_style))
            elif line.startswith('-') or line.startswith('*'):
                story.append(Paragraph(f"• {line[1:].strip()}", normal_style))
            elif line.startswith('**') and line.endswith('**'):
                story.append(Paragraph(f"<b>{line[2:-2]}</b>", normal_style))
            else:
                if len(line) > 20:
                    story.append(Paragraph(line, normal_style))

            # Insert charts at strategic breakpoints
            total_lines = len([l for l in lines if l.strip()])
            if total_charts > 0:
                chart_positions = [
                    int(total_lines * 0.25),
                    int(total_lines * 0.5),
                    int(total_lines * 0.75),
                ]

                if i in chart_positions and chart_inserted < total_charts:
                    charts_to_insert = min(2, total_charts - chart_inserted)
                    for j in range(charts_to_insert):
                        if chart_inserted < total_charts:
                            try:
                                chart_data = base64.b64decode(chart_images[chart_inserted])
                                chart_buffer = io.BytesIO(chart_data)
                                chart_img = Image(chart_buffer, width=5.5 * inch, height=3.5 * inch)
                                story.append(Spacer(1, 10))
                                story.append(chart_img)
                                story.append(Spacer(1, 10))
                                chart_inserted += 1
                            except Exception as e:
                                print(f"Chart insert error {chart_inserted}: {e}")
                                chart_inserted += 1

        if table_pending:
            insert_table(table_pending)

        while chart_inserted < total_charts:
            try:
                chart_data = base64.b64decode(chart_images[chart_inserted])
                chart_buffer = io.BytesIO(chart_data)
                chart_img = Image(chart_buffer, width=5.5 * inch, height=3.5 * inch)
                story.append(Spacer(1, 10))
                story.append(chart_img)
                story.append(Spacer(1, 10))
                chart_inserted += 1
            except Exception as e:
                print(f"Chart insert error {chart_inserted}: {e}")
                chart_inserted += 1

        try:
            doc.build(story)
            return temp_path
        except Exception as e:
            return f"Error creating PDF: {str(e)}"


def generate_report(files, user_prompt, report_title, report_size):
    """Main function to generate report"""

    if not OPENAI_API_KEY:
        return "OpenAI API key not found in Colab secrets. Please add OPENAI_API_KEY to your secrets.", None, None

    if not files:
        return "Please upload at least one file", None, None

    if not user_prompt:
        return "Please provide analysis requirements", None, None

    try:
        # Initialize report generator
        generator = ReportGenerator(OPENAI_API_KEY)

        # Extract content from files
        file_paths = [file.name for file in files]
        content = generator.extract_text_from_files(file_paths)

        # Generate charts
        chart_images = generator.generate_charts_from_data(file_paths)

        # Generate analysis using OpenAI with selected report size
        report_text = generator.analyze_with_openai(content, user_prompt, report_size)

        # Create PDF report
        pdf_path = generator.create_pdf_report(report_text, chart_images, report_title or "Analysis Report")

        return report_text, pdf_path, f"✅ {report_size.title()} report generated successfully! ({len(report_text.split())} words)"

    except Exception as e:
        return f"Error generating report: {str(e)}", None, None

# Create Gradio interface
def create_interface():
    with gr.Blocks(title="AI Report Generator", theme=gr.themes.Soft()) as interface:

        gr.Markdown("""
        # 🤖 AI-Powered Report Generator

        Upload your files (documents, spreadsheets, presentations) and get comprehensive analysis reports with charts and visualizations.

        **Supported formats:** PDF, DOCX, XLSX, CSV, PPTX, TXT
        """)

        with gr.Row():
            with gr.Column(scale=1):
                # Input section
                gr.Markdown("### 📁 Upload Files")
                files = gr.File(
                    label="Select Files",
                    file_count="multiple",
                    file_types=[".pdf", ".docx", ".doc", ".xlsx", ".xls", ".csv", ".pptx", ".ppt", ".txt"]
                )

                gr.Markdown("### 🎯 Analysis Requirements")
                user_prompt = gr.Textbox(
                    label="What would you like to analyze?",
                    placeholder="e.g., 'Analyze sales trends and provide recommendations for Q4 strategy'",
                    lines=3
                )

                report_size = gr.Radio(
                    choices=["brief", "standard", "comprehensive", "detailed"],
                    value="standard",
                    label="📏 Report Size",
                    info="Brief: 1-2 pages | Standard: 2-3 pages | Comprehensive: 3-4 pages | Detailed: 4-5 pages"
                )

                report_title = gr.Textbox(
                    label="Report Title (Optional)",
                    placeholder="e.g., 'Q3 Sales Analysis Report'",
                    value="Analysis Report"
                )

                # Show API key status
                if OPENAI_API_KEY:
                    gr.Markdown("✅ **OpenAI API Key loaded from Colab secrets**")
                else:
                    gr.Markdown("❌ **OpenAI API Key not found in secrets**")

                generate_btn = gr.Button("🚀 Generate Report", variant="primary", size="lg")

            with gr.Column(scale=2):
                # Output section
                gr.Markdown("### 📊 Generated Report")

                status = gr.Textbox(label="Status", interactive=False)

                report_preview = gr.Markdown(
                    label="📊 Report Preview",
                    value="Generated report will appear here...",
                    height=600
                )

                pdf_download = gr.File(
                    label="📄 Download PDF Report",
                    interactive=False
                )

        # Event handlers
        generate_btn.click(
            fn=generate_report,
            inputs=[files, user_prompt, report_title, report_size],
            outputs=[report_preview, pdf_download, status]
        )

        gr.Markdown("""
        ### 💡 Report Size Guide:

        📋 **Brief (1-2 pages)**: Executive summary with key insights and recommendations
        - 4 main sections, ~800-1200 words
        - Focus on critical findings and action items

        📊 **Standard (2-3 pages)**: Balanced analysis with detailed insights
        - 6 main sections, ~1800-2800 words
        - Comprehensive findings with strategic recommendations

        📈 **Comprehensive (3-4 pages)**: Thorough investigation with statistical analysis
        - 8+ sections with sub-sections, ~2800-4200 words
        - Deep dive analysis with risk assessment and implementation roadmap

        📋 **Detailed (4-5 pages)**: Exhaustive research-level analysis
        - 10+ sections with extensive sub-sections, ~4000-6000 words
        - Multi-dimensional analysis with advanced statistical insights

        ### 🎯 Example Prompts by Size:

        **Brief**: *"Quick overview of sales performance with key recommendations"*
        **Standard**: *"Analyze customer satisfaction survey and provide actionable insights"*
        **Comprehensive**: *"Detailed financial analysis with risk assessment and strategic planning"*
        **Detailed**: *"Exhaustive market research with competitive analysis and implementation blueprint"*
        """)



    return interface

if __name__ == "__main__":
    # Create and launch the interface
    app = create_interface()
    app.launch(
        share=True,
        debug=True,
        show_api=False
    )
